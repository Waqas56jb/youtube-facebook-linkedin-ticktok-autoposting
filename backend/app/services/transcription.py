from typing import List, Optional, Tuple
import os
import tempfile
import shutil
import mimetypes
import chardet
import csv

from pypdf import PdfReader
from docx import Document as DocxDocument
from bs4 import BeautifulSoup
from markdown_it import MarkdownIt
from pptx import Presentation

try:
    import pypandoc  # optional
except Exception:
    pypandoc = None  # type: ignore

from fastapi import UploadFile, HTTPException

from backend.app.services.whisper import get_whisper_model
from backend.app.models import TranscriptSegment


def save_upload_to_temp(upload: UploadFile) -> str:
    suffix = os.path.splitext(upload.filename or "uploaded")[1]
    tmp_fd, tmp_path = tempfile.mkstemp(suffix=suffix)
    os.close(tmp_fd)
    with open(tmp_path, "wb") as out_f:
        shutil.copyfileobj(upload.file, out_f)
    return tmp_path


def detect_mime_type(file_path: str, fallback_name: Optional[str]) -> str:
    guessed, _ = mimetypes.guess_type(fallback_name or file_path)
    return guessed or "application/octet-stream"


def extract_text_from_pdf(file_path: str) -> str:
    reader = PdfReader(file_path)
    pages_text: List[str] = []
    for page in reader.pages:
        try:
            pages_text.append(page.extract_text() or "")
        except Exception:
            pages_text.append("")
    return "\n\n".join([t.strip() for t in pages_text if t and t.strip()])


def extract_text_from_docx(file_path: str) -> str:
    doc = DocxDocument(file_path)
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text and p.text.strip()]
    return "\n\n".join(paragraphs)


def extract_text_from_txt(file_path: str) -> str:
    with open(file_path, "rb") as f:
        raw = f.read()
    detected = chardet.detect(raw)
    encoding = (detected.get("encoding") or "utf-8")
    try:
        return raw.decode(encoding, errors="replace")
    except Exception:
        return raw.decode("utf-8", errors="replace")


def extract_text_from_html(file_path: str) -> str:
    html_text = extract_text_from_txt(file_path)
    soup = BeautifulSoup(html_text, "html.parser")
    return soup.get_text("\n", strip=True)


def extract_text_from_markdown(file_path: str) -> str:
    md_src = extract_text_from_txt(file_path)
    try:
        html = MarkdownIt().render(md_src)
        soup = BeautifulSoup(html, "html.parser")
        return soup.get_text("\n", strip=True)
    except Exception:
        return md_src


def extract_text_from_rtf(file_path: str) -> str:
    if pypandoc is not None:
        try:
            return pypandoc.convert_file(file_path, "plain")
        except Exception:
            pass
    return extract_text_from_txt(file_path)


def extract_text_from_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    texts: List[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
    return "\n\n".join([t.strip() for t in texts if t and t.strip()])


def extract_text_from_csv(file_path: str) -> str:
    with open(file_path, "rb") as f:
        raw = f.read()
    detected = chardet.detect(raw)
    encoding = (detected.get("encoding") or "utf-8")
    lines: List[str] = []
    try:
        with open(file_path, newline="", encoding=encoding, errors="replace") as csvfile:
            reader = csv.reader(csvfile)
            for row in reader:
                lines.append(", ".join([col.strip() for col in row]))
    except Exception:
        return extract_text_from_txt(file_path)
    return "\n".join(lines)


def transcribe_media(file_path: str):
    model = get_whisper_model()
    _verify_media_readable(file_path)
    from backend.app.config import settings

    def _run(path: str, force_lang: bool = False):
        try:
            return model.transcribe(
                path,
                beam_size=getattr(settings, 'whisper_beam_size', 1),
                vad_filter=True,
                vad_parameters={"min_silence_duration_ms": 250},
                word_timestamps=False,
                condition_on_previous_text=False,
                chunk_length=getattr(settings, 'whisper_chunk_length', 30),
                language=(getattr(settings, 'whisper_language', 'en') or 'en') if force_lang else (getattr(settings, 'whisper_language', None) or None),
                task="transcribe",
                temperature=0.0,
            )
        except Exception as _:
            # Second chance with forced English
            return model.transcribe(
                path,
                beam_size=1,
                vad_filter=True,
                vad_parameters={"min_silence_duration_ms": 250},
                word_timestamps=False,
                condition_on_previous_text=False,
                chunk_length=30,
                language='en' if force_lang else None,
                task="transcribe",
                temperature=0.0,
            )

    # First attempt on original media
    segments_iter, info = _run(file_path, force_lang=False)
    segments: List[TranscriptSegment] = []
    texts: List[str] = []
    for seg in segments_iter:
        text_clean = (seg.text or "").strip()
        segments.append(TranscriptSegment(start=seg.start, end=seg.end, text=text_clean))
        if text_clean:
            texts.append(text_clean)
    full_text = " ".join(texts).strip()
    if full_text:
        return full_text, segments, getattr(info, "language", None), getattr(info, "duration", None)

    # Fallback: extract audio to 16kHz mono WAV and retry
    try:
        import subprocess, tempfile
        with tempfile.NamedTemporaryFile(suffix='.wav', delete=True) as tmp:
            wav_path = tmp.name
            cmd = [
                'ffmpeg','-y','-i', file_path,
                '-vn','-acodec','pcm_s16le','-ar','16000','-ac','1', wav_path
            ]
            subprocess.run(cmd, capture_output=True, check=True)
            segments_iter2, info2 = _run(wav_path, force_lang=True)
            segments2: List[TranscriptSegment] = []
            texts2: List[str] = []
            for seg in segments_iter2:
                text_clean = (seg.text or "").strip()
                segments2.append(TranscriptSegment(start=seg.start, end=seg.end, text=text_clean))
                if text_clean:
                    texts2.append(text_clean)
            full_text2 = " ".join(texts2).strip()
            if full_text2:
                return full_text2, segments2, getattr(info2, "language", None), getattr(info2, "duration", None)
    except Exception:
        pass

    # Return empty but valid structure if all attempts fail
    return "", [], None, None

def _verify_media_readable(src_path: str) -> None:
    """Lightweight check that ffmpeg/ffprobe can read the file without creating new artifacts."""
    try:
        import subprocess, shlex
        cmd = f"ffprobe -v error -show_format -show_streams -of json {shlex.quote(src_path)}"
        subprocess.run(cmd, shell=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, check=True)
    except Exception:
        # Don't block; whisper will attempt decode. This avoids creating extra files per user's request.
        pass


DOC_EXTS = (
    ".pdf", ".docx", ".txt", ".rtf", ".html", ".htm", ".md", ".markdown", ".pptx", ".csv", ".srt", ".vtt"
)


def is_document_file(name_lower: str, mime: str) -> bool:
    return any(name_lower.endswith(ext) for ext in DOC_EXTS) or (
        mime.startswith("application/pdf")
        or mime.endswith("msword")
        or "officedocument" in mime
        or mime.startswith("text/")
        or mime in (
            "application/rtf",
            "text/rtf",
            "text/html",
            "text/markdown",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "text/csv",
            "application/csv",
            "text/vtt",
            "application/x-subrip",
        )
    )


def extract_document_text(temp_path: str, name_lower: str, mime: str) -> str:
    if name_lower.endswith(".pdf") or mime.startswith("application/pdf"):
        return extract_text_from_pdf(temp_path)
    if name_lower.endswith(".docx") or "officedocument.wordprocessingml.document" in mime:
        return extract_text_from_docx(temp_path)
    if name_lower.endswith(".rtf") or mime in ("application/rtf", "text/rtf"):
        return extract_text_from_rtf(temp_path)
    if name_lower.endswith(".html") or name_lower.endswith(".htm") or mime == "text/html":
        return extract_text_from_html(temp_path)
    if name_lower.endswith(".md") or name_lower.endswith(".markdown") or mime == "text/markdown":
        return extract_text_from_markdown(temp_path)
    if name_lower.endswith(".pptx") or mime == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        return extract_text_from_pptx(temp_path)
    if name_lower.endswith(".csv") or mime in ("text/csv", "application/csv"):
        return extract_text_from_csv(temp_path)
    # subtitles and any other text-like
    return extract_text_from_txt(temp_path)


